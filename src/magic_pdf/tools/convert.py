"""Format conversion tools (to and from PDF)."""

import os
import shutil
import subprocess
import tempfile

import fitz  # PyMuPDF
from PIL import Image

from .utils import (
    validate_file_exists,
    validate_file_size,
    validate_extension,
    generate_output_path,
    ensure_output_dir,
    logger,
)


def _find_libreoffice() -> str | None:
    """Find LibreOffice executable."""
    for name in ["libreoffice", "soffice"]:
        path = shutil.which(name)
        if path:
            return path
    # Windows common paths
    for path in [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]:
        if os.path.isfile(path):
            return path
    return None


def _convert_with_libreoffice(input_path: str, output_dir: str) -> str | None:
    """Convert a document to PDF using LibreOffice. Returns output path or None."""
    lo = _find_libreoffice()
    if not lo:
        return None

    try:
        subprocess.run(
            [lo, "--headless", "--convert-to", "pdf", "--outdir", output_dir, input_path],
            check=True,
            capture_output=True,
            timeout=120,
        )
        base = os.path.splitext(os.path.basename(input_path))[0]
        output_path = os.path.join(output_dir, f"{base}.pdf")
        if os.path.isfile(output_path):
            return output_path
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired, FileNotFoundError):
        pass
    return None


def docx_to_pdf(file_path: str, output_path: str | None = None) -> str:
    """Convert a Word document (.docx) to PDF.

    Uses LibreOffice when available for high-fidelity conversion,
    falls back to pure Python (python-docx + reportlab) otherwise.

    Args:
        file_path: Absolute path to the .docx file
        output_path: Optional output path (default: same directory with .pdf extension)
    """
    validate_file_exists(file_path)
    validate_file_size(file_path)
    validate_extension(file_path, [".docx"])

    if output_path is None:
        output_path = generate_output_path(file_path, "converted", ".pdf")
    ensure_output_dir(output_path)

    # Try LibreOffice first
    output_dir = os.path.dirname(output_path) or "."
    lo_result = _convert_with_libreoffice(file_path, output_dir)
    if lo_result:
        if lo_result != output_path:
            os.rename(lo_result, output_path)
        return f"Converted (LibreOffice): {output_path}"

    # Fallback: python-docx + reportlab
    logger.warning(
        "LibreOffice not found. Using pure Python fallback — fidelity may be reduced."
    )
    from docx import Document
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet

    doc = Document(file_path)
    pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
    styles = getSampleStyleSheet()
    story = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            # Map heading styles
            if para.style and para.style.name.startswith("Heading"):
                try:
                    level = int(para.style.name.replace("Heading", "").strip())
                    style = styles.get(f"Heading{min(level, 6)}", styles["Normal"])
                except ValueError:
                    style = styles["Normal"]
            else:
                style = styles["Normal"]
            story.append(Paragraph(text, style))
            story.append(Spacer(1, 6))

    if not story:
        story.append(Paragraph(" ", styles["Normal"]))

    pdf_doc.build(story)
    return f"Converted (fallback): {output_path}"


def image_to_pdf(
    file_paths: list[str],
    output_path: str,
) -> str:
    """Convert image(s) to PDF. Supported formats: JPEG, PNG, BMP, TIFF, GIF.

    Args:
        file_paths: List of absolute paths to image files
        output_path: Absolute path for the output PDF file
    """
    if not file_paths:
        raise ValueError("No image files provided")

    allowed = [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif", ".gif"]

    images = []
    for fp in file_paths:
        validate_file_exists(fp)
        validate_file_size(fp)
        validate_extension(fp, allowed)
        img = Image.open(fp)
        if img.mode == "RGBA":
            img = img.convert("RGB")
        elif img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        images.append(img)

    ensure_output_dir(output_path)

    if len(images) == 1:
        images[0].save(output_path, "PDF")
    else:
        images[0].save(output_path, "PDF", save_all=True, append_images=images[1:])

    for img in images:
        img.close()

    return f"Images converted to PDF: {output_path}"


def html_to_pdf(
    source: str,
    output_path: str,
) -> str:
    """Convert HTML to PDF. Accepts either a raw HTML string or a URL to fetch and convert.

    Args:
        source: Raw HTML string or URL (http:// or https://)
        output_path: Absolute path for the output PDF file
    """
    ensure_output_dir(output_path)

    # Determine if source is a URL or HTML string
    if source.strip().startswith(("http://", "https://")):
        import urllib.request

        try:
            with urllib.request.urlopen(source, timeout=30) as response:
                html_content = response.read().decode("utf-8", errors="replace")
        except Exception as e:
            raise RuntimeError(f"Failed to fetch URL: {e}")
    else:
        html_content = source

    from xhtml2pdf import pisa

    with open(output_path, "wb") as f:
        result = pisa.CreatePDF(html_content, dest=f)
        if result.err:
            raise RuntimeError(f"HTML to PDF conversion failed with {result.err} errors")

    return f"HTML converted to PDF: {output_path}"


def excel_to_pdf(file_path: str, output_path: str | None = None) -> str:
    """Convert an Excel spreadsheet (.xlsx/.xls) to PDF.

    Uses LibreOffice when available, falls back to openpyxl + reportlab.

    Args:
        file_path: Absolute path to the Excel file
        output_path: Optional output path (default: same directory with .pdf extension)
    """
    validate_file_exists(file_path)
    validate_file_size(file_path)
    validate_extension(file_path, [".xlsx", ".xls"])

    if output_path is None:
        output_path = generate_output_path(file_path, "converted", ".pdf")
    ensure_output_dir(output_path)

    # Try LibreOffice first
    output_dir = os.path.dirname(output_path) or "."
    lo_result = _convert_with_libreoffice(file_path, output_dir)
    if lo_result:
        if lo_result != output_path:
            os.rename(lo_result, output_path)
        return f"Converted (LibreOffice): {output_path}"

    # Fallback: openpyxl + reportlab
    logger.warning(
        "LibreOffice not found. Using pure Python fallback — fidelity may be reduced."
    )
    from openpyxl import load_workbook
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
    from reportlab.lib import colors

    wb = load_workbook(file_path, data_only=True)
    ws = wb.active

    data = []
    for row in ws.iter_rows(values_only=True):
        data.append([str(cell) if cell is not None else "" for cell in row])

    if not data:
        data = [["(empty spreadsheet)"]]

    pdf_doc = SimpleDocTemplate(output_path, pagesize=landscape(A4))
    table = Table(data)
    table.setStyle(
        TableStyle(
            [
                ("GRID", (0, 0), (-1, -1), 0.5, colors.black),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]
        )
    )
    pdf_doc.build([table])
    return f"Converted (fallback): {output_path}"


def powerpoint_to_pdf(file_path: str, output_path: str | None = None) -> str:
    """Convert a PowerPoint presentation (.pptx/.ppt) to PDF.

    Uses LibreOffice when available, falls back to python-pptx + reportlab.

    Args:
        file_path: Absolute path to the PowerPoint file
        output_path: Optional output path (default: same directory with .pdf extension)
    """
    validate_file_exists(file_path)
    validate_file_size(file_path)
    validate_extension(file_path, [".pptx", ".ppt"])

    if output_path is None:
        output_path = generate_output_path(file_path, "converted", ".pdf")
    ensure_output_dir(output_path)

    # Try LibreOffice first
    output_dir = os.path.dirname(output_path) or "."
    lo_result = _convert_with_libreoffice(file_path, output_dir)
    if lo_result:
        if lo_result != output_path:
            os.rename(lo_result, output_path)
        return f"Converted (LibreOffice): {output_path}"

    # Fallback: python-pptx + reportlab
    logger.warning(
        "LibreOffice not found. Using pure Python fallback — fidelity may be reduced."
    )
    from pptx import Presentation
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.pdfgen import canvas as pdf_canvas

    prs = Presentation(file_path)
    c = pdf_canvas.Canvas(output_path, pagesize=landscape(A4))
    width, height = landscape(A4)

    for slide_num, slide in enumerate(prs.slides, 1):
        y = height - 72
        c.setFont("Helvetica-Bold", 10)
        c.drawString(72, y, f"Slide {slide_num}")
        y -= 24
        c.setFont("Helvetica", 10)

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text
                    if text.strip():
                        # Wrap long lines
                        while len(text) > 100:
                            c.drawString(72, y, text[:100])
                            text = text[100:]
                            y -= 14
                            if y < 72:
                                c.showPage()
                                y = height - 72
                        c.drawString(72, y, text)
                        y -= 14
                        if y < 72:
                            c.showPage()
                            y = height - 72
        c.showPage()

    c.save()
    return f"Converted (fallback): {output_path}"


def pdf_to_image(
    file_path: str,
    output_dir: str | None = None,
    output_format: str = "png",
    dpi: int = 150,
    pages: str | None = None,
) -> str:
    """Convert PDF page(s) to images.

    Args:
        file_path: Absolute path to the PDF file
        output_dir: Directory for output images (default: same as input file)
        output_format: Output format - 'png' or 'jpeg' (default: png)
        dpi: Resolution in DPI (default: 150)
        pages: Page specification - single page '3', range '1-5', or 'all' (default: all)
    """
    from .utils import validate_pdf_file

    validate_pdf_file(file_path)
    validate_file_size(file_path)

    if output_format.lower() not in ("png", "jpeg", "jpg"):
        raise ValueError(f"Unsupported format: {output_format}. Use 'png' or 'jpeg'.")

    fmt = "png" if output_format.lower() == "png" else "jpeg"

    if output_dir is None:
        output_dir = os.path.dirname(file_path)
    os.makedirs(output_dir, exist_ok=True)

    doc = fitz.open(file_path)
    total_pages = len(doc)

    # Parse page specification
    if pages is None or pages.lower() == "all":
        page_nums = list(range(total_pages))
    elif "-" in pages:
        start, end = pages.split("-", 1)
        start_idx = int(start) - 1
        end_idx = int(end)
        if start_idx < 0 or end_idx > total_pages or start_idx >= end_idx:
            raise ValueError(f"Invalid page range: {pages} (PDF has {total_pages} pages)")
        page_nums = list(range(start_idx, end_idx))
    else:
        page_num = int(pages) - 1
        if page_num < 0 or page_num >= total_pages:
            raise ValueError(f"Invalid page: {pages} (PDF has {total_pages} pages)")
        page_nums = [page_num]

    base_name = os.path.splitext(os.path.basename(file_path))[0]
    output_files = []
    mat = fitz.Matrix(dpi / 72, dpi / 72)

    for page_num in page_nums:
        page = doc[page_num]
        pix = page.get_pixmap(matrix=mat)

        ext = "png" if fmt == "png" else "jpg"
        output_file = os.path.join(output_dir, f"{base_name}_page{page_num + 1}.{ext}")

        if fmt == "png":
            pix.save(output_file)
        else:
            # Convert via Pillow for JPEG
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            img.save(output_file, format="JPEG", quality=90)
            img.close()

        output_files.append(output_file)

    doc.close()
    return f"Converted {len(output_files)} page(s) to {fmt}: {', '.join(output_files)}"


def register(mcp):
    """Register format conversion tools with the MCP server."""
    mcp.tool()(docx_to_pdf)
    mcp.tool()(image_to_pdf)
    mcp.tool()(html_to_pdf)
    mcp.tool()(excel_to_pdf)
    mcp.tool()(powerpoint_to_pdf)
    mcp.tool()(pdf_to_image)
