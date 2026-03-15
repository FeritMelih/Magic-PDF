"""Shared test fixtures for Magic-PDF tests."""

import os
import pytest
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas as pdf_canvas


@pytest.fixture
def sample_pdf(tmp_path):
    """Create a simple 3-page PDF for testing."""
    path = str(tmp_path / "sample.pdf")
    c = pdf_canvas.Canvas(path, pagesize=A4)

    c.setFont("Helvetica", 12)
    c.drawString(72, 700, "Hello World - Page 1")
    c.drawString(72, 680, "The quick brown fox jumps over the lazy dog")
    c.showPage()

    c.setFont("Helvetica", 12)
    c.drawString(72, 700, "Hello World - Page 2")
    c.drawString(72, 680, "Python is a great programming language")
    c.showPage()

    c.setFont("Helvetica", 12)
    c.drawString(72, 700, "Hello World - Page 3")
    c.drawString(72, 680, "PDF manipulation is useful")
    c.showPage()

    c.save()
    return path


@pytest.fixture
def sample_pdf_searchable(tmp_path):
    """Create a PDF with specific searchable text."""
    path = str(tmp_path / "searchable.pdf")
    c = pdf_canvas.Canvas(path, pagesize=A4)

    c.setFont("Helvetica", 12)
    c.drawString(72, 700, "The quick brown fox jumps over the lazy dog")
    c.drawString(72, 680, "Python programming is fun and productive")
    c.drawString(72, 660, "Machine learning uses Python extensively")
    c.showPage()

    c.setFont("Helvetica", 12)
    c.drawString(72, 700, "Another page with Python references")
    c.drawString(72, 680, "Data science and Python go hand in hand")
    c.showPage()

    c.save()
    return path


@pytest.fixture
def sample_pdf_pair(tmp_path):
    """Create two small PDFs for merge testing."""
    paths = []
    for i in range(1, 3):
        path = str(tmp_path / f"doc{i}.pdf")
        c = pdf_canvas.Canvas(path, pagesize=A4)
        c.setFont("Helvetica", 12)
        c.drawString(72, 700, f"Document {i} - Page 1")
        c.showPage()
        c.drawString(72, 700, f"Document {i} - Page 2")
        c.showPage()
        c.save()
        paths.append(path)
    return paths


@pytest.fixture
def sample_image(tmp_path):
    """Create a simple test image."""
    from PIL import Image

    path = str(tmp_path / "test_image.png")
    img = Image.new("RGB", (200, 300), color=(255, 0, 0))
    img.save(path)
    img.close()
    return path


@pytest.fixture
def sample_docx(tmp_path):
    """Create a simple .docx file for testing."""
    from docx import Document

    path = str(tmp_path / "test.docx")
    doc = Document()
    doc.add_heading("Test Document", level=1)
    doc.add_paragraph("This is a test paragraph.")
    doc.add_paragraph("Another paragraph with some content.")
    doc.save(path)
    return path


@pytest.fixture
def sample_xlsx(tmp_path):
    """Create a simple .xlsx file for testing."""
    from openpyxl import Workbook

    path = str(tmp_path / "test.xlsx")
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Name"
    ws["B1"] = "Value"
    ws["A2"] = "Alpha"
    ws["B2"] = 100
    ws["A3"] = "Beta"
    ws["B3"] = 200
    wb.save(path)
    return path


@pytest.fixture
def sample_pptx(tmp_path):
    """Create a simple .pptx file for testing."""
    from pptx import Presentation

    path = str(tmp_path / "test.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Slide"
    slide.placeholders[1].text = "Test content"
    prs.save(path)
    return path


@pytest.fixture
def sample_form_pdf(tmp_path):
    """Create a PDF with form fields using reportlab."""
    path = str(tmp_path / "form.pdf")
    c = pdf_canvas.Canvas(path, pagesize=A4)

    c.setFont("Helvetica", 12)
    c.drawString(72, 750, "Test Form")

    form = c.acroForm
    form.textfield(
        name="name",
        x=72,
        y=700,
        width=200,
        height=20,
        borderWidth=1,
        fieldFlags="",
    )
    form.textfield(
        name="email",
        x=72,
        y=660,
        width=200,
        height=20,
        borderWidth=1,
        fieldFlags="",
    )

    c.save()
    return path
